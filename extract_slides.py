#!/usr/bin/env python3
"""
Extract PowerPoint slides as images
"""
import os
from pptx import Presentation
from PIL import Image
import io

def extract_slides_to_images(pptx_path, output_dir):
    """
    Extract all slides from a PowerPoint presentation as images.
    Note: This extracts embedded images from slides, not full slide renders.
    """
    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Load the presentation
    print(f"Loading presentation: {pptx_path}")
    prs = Presentation(pptx_path)
    
    # Track all images found
    image_count = 0
    slide_images = {}  # Track which images belong to which slides
    
    # Iterate through all slides
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"\nProcessing slide {slide_idx}...")
        slide_images[slide_idx] = []
        
        # Extract images from shapes
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                image_count += 1
                # Get image data
                image = shape.image
                image_bytes = image.blob
                
                # Save the image
                ext = image.ext
                filename = f"slide{slide_idx}_image{shape_idx}.{ext}"
                filepath = os.path.join(output_dir, filename)
                
                with open(filepath, "wb") as f:
                    f.write(image_bytes)
                
                print(f"  Extracted image: {filename}")
                slide_images[slide_idx].append(filename)
            
            # Check if shape contains a picture
            if shape.shape_type == 13:  # Picture shape type
                try:
                    if hasattr(shape, 'image'):
                        continue  # Already handled above
                except:
                    pass
    
    print(f"\nTotal images extracted: {image_count}")
    
    # Report which slides had images
    print("\nImages by slide:")
    for slide_num, images in slide_images.items():
        if images:
            print(f"  Slide {slide_num}: {len(images)} image(s)")
    
    return image_count

if __name__ == "__main__":
    pptx_file = "/Users/proth/Downloads/NIL GAMEPLAN Funding Deck 7-29A.pptx"
    output_folder = "/Users/proth/repos/nil-gameplan/slides/"
    
    try:
        count = extract_slides_to_images(pptx_file, output_folder)
        print(f"\nExtraction complete! {count} images saved to {output_folder}")
    except Exception as e:
        print(f"Error: {e}")
        print("\nNote: python-pptx can only extract embedded images, not render full slides.")
        print("For full slide rendering, you would need PowerPoint or LibreOffice installed.")